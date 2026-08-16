# -*- coding: utf-8 -*-
"""CI 校验：markdown fence 偶数 + 敏感信息扫描 + PPTX 越界检查
用法: python ci_validate.py [--skip-pptx]
退出码: 0 = 全部通过; 1 = 有问题
"""
import glob
import os
import re
import sys

# 仓库根 = 本文件(ci/ci_validate.py)的上两级目录
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
COURSE = ROOT
SKIP_PPTX = "--skip-pptx" in sys.argv

# 与 .gitignore 保持一致：这些文件不入库，不参与校验（本地可能存在）
IGNORED = {
    os.path.join(COURSE, "tools", "r.ps1"),
    os.path.join(COURSE, "misc", "00-Kubernetes 核心知识.md"),
}

issues = []


def check_fence():
    """所有 .md 代码围栏必须偶数"""
    files = glob.glob(os.path.join(COURSE, "**", "*.md"), recursive=True)
    for f in files:
        if f in IGNORED:
            continue
        with open(f, encoding="utf-8") as fh:
            c = fh.read()
        n = len(re.findall(r"(?m)^```", c))
        if n % 2 != 0:
            issues.append(f"fence 奇数({n}): {os.path.relpath(f, COURSE)}")
    print(f"[1/3] fence 检查: {len(files)} 个 md 文件, 问题 {len([i for i in issues if 'fence' in i])}")


def check_secrets():
    """敏感信息扫描：禁止公网 IP / 密码 / 私钥出现在将入库的文件中"""
    pats = [
        (r"\b1\.95\.127\.232\b|\b101\.245\.100\.229\b|\b1\.95\.74\.210\b", "公网IP"),
        (r"2wsx#EDC", "密码"),
        (r"BEGIN (RSA |EC |OPENSSH )?PRIVATE KEY", "私钥"),
        (r"client-key-data:\s*[A-Za-z0-9+/]{50,}", "客户端私钥数据"),
    ]
    files = glob.glob(os.path.join(COURSE, "**", "*"), recursive=True)
    files = [f for f in files if os.path.isfile(f)
             and not f.endswith(".pptx") and not f.endswith(".gitignore")
             and f not in IGNORED
             and os.path.basename(f) != "ci_validate.py"]  # 排除扫描脚本自身（含模式字符串）
    for f in files:
        try:
            with open(f, encoding="utf-8", errors="ignore") as fh:
                c = fh.read()
        except Exception:
            continue
        for pat, name in pats:
            if re.search(pat, c):
                issues.append(f"敏感信息({name}): {os.path.relpath(f, COURSE)}")
    print(f"[2/3] 敏感信息扫描: 问题 {len([i for i in issues if '敏感' in i])}")


def check_pptx():
    if SKIP_PPTX:
        print("[3/3] PPTX 越界检查: 跳过 (--skip-pptx)")
        return
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        print("[3/3] PPTX 越界检查: 跳过 (无 python-pptx)")
        return
    files = glob.glob(os.path.join(COURSE, "ppt", "**", "*.pptx"), recursive=True)
    n_issue = 0
    for f in files:
        try:
            p = Presentation(f)
        except Exception as e:
            issues.append(f"PPTX 打开失败: {os.path.relpath(f, COURSE)} ({e})")
            continue
        W, H = p.slide_width, p.slide_height
        for i, s in enumerate(p.slides, 1):
            for shp in s.shapes:
                if shp.left is None:
                    continue
                if shp.left < 0 or shp.top < 0 or \
                   shp.left + shp.width > W + Emu(50000) or \
                   shp.top + shp.height > H + Emu(50000):
                    issues.append(f"越界 P{i}: {os.path.relpath(f, COURSE)}")
                    n_issue += 1
                    break
    print(f"[3/3] PPTX 越界检查: {len(files)} 个 pptx, 问题 {n_issue}")


def main():
    check_fence()
    check_secrets()
    check_pptx()
    if issues:
        print("\n❌ 发现问题:")
        for i in issues:
            print("  -", i)
        sys.exit(1)
    print("\n✅ 全部校验通过")


if __name__ == "__main__":
    main()
